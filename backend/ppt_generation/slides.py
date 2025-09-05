from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import pandas as pd
import datetime
from copy import deepcopy


def set_cell_border(cell, border_color="000000", border_width="12700"):
    """
    Add border to a pptx table cell.
    Default: Black, ~0.5pt width.
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for line in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = tcPr.find(qn(line))
        if ln is None:
            ln = parse_xml(
                f'<{line} w="{border_width}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:solidFill><a:srgbClr val="{border_color}"/></a:solidFill></{line}>'
            )
            tcPr.append(ln)


def style_header_cell(cell, font_size=12):
    """
    Apply styling to a header cell.
    """
    cell.text_frame.word_wrap = True  # ✅ enable wrapping
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p.alignment = PP_ALIGN.CENTER

    # Background + border
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
    set_cell_border(cell)


def format_value(value):
    if pd.isna(value):
        return ""
    # Handle Python datetime.datetime or pandas Timestamp objects
    if isinstance(value, (datetime.datetime, datetime.date, pd.Timestamp)):
        # For datetime.datetime convert to date to remove time part
        if isinstance(value, datetime.datetime):
            value = value.date()
        return value.strftime("%d-%m-%Y")  # Format as DD-MM-YYYY
    # If value is string that contains date-time, parse and format
    elif isinstance(value, str):
        try:
            # Try to parse string as date-time with time part, then format
            parsed = datetime.datetime.strptime(value, "%Y-%m-%d %H:%M:%S")
            return parsed.strftime("%d-%m-%Y")
        except Exception:
            # If parsing fails, return original string
            return value
    return str(value)


# def style_data_cell(cell, value, font_size=10, fill_color=None):
#     """
#     Apply styling to a data cell with support for custom background color.
#     fill_color: Excel ARGB color string such as 'FFB6DDE8'
#     """
#     cell.text_frame.text = format_value(value)
#     cell.text_frame.word_wrap = True  # ✅ enable wrapping

#     for p in cell.text_frame.paragraphs:
#         for run in p.runs:
#             run.font.size = Pt(font_size)
#             run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
#         p.alignment = PP_ALIGN.LEFT

#     # Background + border
#     cell.fill.solid()
#     if fill_color and isinstance(fill_color, str) and len(fill_color) == 8:
#         # Remove alpha ('FF'), use RRGGBB for pptx
#         rgb = fill_color[-6:]
#         cell.fill.fore_color.rgb = RGBColor.from_string(rgb)
#     else:
#         cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # default White

#     set_cell_border(cell)
def style_data_cell(cell, cell_info, font_size=10):
    """
    Apply styling to a data cell.
    Handles both dict {value, style} and plain strings/numbers.
    Ensures text always shows.
    """
    # --- Extract info ---
    if isinstance(cell_info, dict):
        value = cell_info.get("value", "")
        style = cell_info.get("style") or {}
    else:
        value = cell_info
        style = {}

    # Always convert to safe string
    value = format_value(value) if value is not None else ""
    # print("Placing value in PPT:", repr(value))  # ✅ Debug log

    fill_color = style.get("fill_color") or style.get("fill")
    font_color = style.get("font_color")
    bold = style.get("font_bold", False)
    italic = style.get("font_italic", False)

    # --- Always reset + add text ---
    text_frame = cell.text_frame
    text_frame.clear()

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = value  # ✅ Safe assignment as string
    text_frame.word_wrap = True

    # --- Font styles ---
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic

    if font_color and isinstance(font_color, str):
        rgb = font_color.replace("#", "")
        if len(rgb) == 8:  # strip ARGB
            rgb = rgb[-6:]
        try:
            run.font.color.rgb = RGBColor.from_string(rgb)
        except Exception:
            run.font.color.rgb = RGBColor(0, 0, 0)
    else:
        run.font.color.rgb = RGBColor(0, 0, 0)

    # --- Alignment ---
    p.alignment = PP_ALIGN.LEFT

    # --- Background fill ---
    cell.fill.solid()
    if fill_color and isinstance(fill_color, str):
        rgb = fill_color.replace("#", "")
        if len(rgb) == 8:
            rgb = rgb[-6:]
        try:
            cell.fill.fore_color.rgb = RGBColor.from_string(rgb)
        except Exception:
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    else:
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # --- Borders ---
    set_cell_border(cell)

def adjust_row_height(table, row_idx, df_row, col_widths, base_height=0.4, max_height=1.5):
    max_lines = 1
    for col_idx, value in enumerate(df_row):
        if col_idx >= len(col_widths): continue
        text = format_value(value)
        approx_chars_per_line = int(col_widths[col_idx].inches * 12)
        line_count = (len(text) // approx_chars_per_line) + 1
        if line_count > max_lines:
            max_lines = line_count
    height_in_inches = min(base_height * max_lines, max_height)
    table.rows[row_idx].height = Inches(height_in_inches)

def insert_slide(prs, layout, index):
    """
    Create a slide using `layout` and move it to `index` (0-based).
    This correctly moves the *newly added* slide's sldId node.
    """
    # add at the end first
    new_slide = prs.slides.add_slide(layout)

    sldIdLst = prs.slides._sldIdLst               # XML <p:sldIdLst>
    new_sldId = sldIdLst[-1]                       # the sldId for the slide just added
    sldIdLst.remove(new_sldId)

    # clamp index to valid bounds
    index = max(0, min(index, len(sldIdLst)))
    sldIdLst.insert(index, new_sldId)

    return new_slide

# def create_ppt_from_excel(template_path, excel_data, output_path, rows_per_slide=5):
#     """
#     Creates a PPT from Excel sheet data (dict of {sheet_name: dataframe}).
#     Keeps the first slide of template as-is, starts inserting data from slide 2.
#     """
#     prs = Presentation(template_path)

#     for sheet_name, df in excel_data.items():

#         if df.empty or len(df.columns) == 0:
#             continue

#         total_rows = len(df)
#         start_row = 0
#         slide_count = 1

#         while start_row < total_rows:
#             slide_layout = prs.slide_layouts[1]
#             slide = prs.slides.add_slide(slide_layout)
#             # slide_idx = len(prs.slides) - 1
#             # if slide_idx > 2:
#             #       prs.slides._sldIdLst.insert(2, prs.slides._sldIdLst[-1])
#             #       del prs.slides._sldIdLst[-1]
#             #       slide = prs.slides[2]  # Now the new slide is in the correct position
#             # Move slide before last two template slides:
#             slide = add_slide_at_index(prs, slide_layout, 2 + (slide_count - 1))

#             # Title
#             title_text = f"{sheet_name} (Contd..)" if slide_count > 1 else sheet_name
#             if slide.shapes.title:
#                 slide.shapes.title.text = title_text
#             else:
#                 left, top, width, height = Inches(0.5), Inches(0.3), Inches(9), Inches(1)
#                 textbox = slide.shapes.add_textbox(left, top, width, height)
#                 textbox.text_frame.text = title_text

#             # Table
#             rows = min(rows_per_slide, total_rows - start_row) + 1
#             cols = len(df.columns)
#             left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.8 * rows)

#             table = slide.shapes.add_table(rows, cols, left, top, width, height).table

#             # Column widths (customize as needed)
#             total_table_width = Inches(12)  # max table width on slide
#             num_cols = len(df.columns)
#             # default_col_width = (total_table_width / num_cols)
#             col_widths = [Inches(total_table_width.inches / num_cols)] * num_cols

#             for i, w in enumerate(col_widths[:cols]):
#                 table.columns[i].width = int(w.emu)

#             # Header row
#             for col_idx, col_name in enumerate(df.columns):
#                 cell = table.cell(0, col_idx)
#                 cell.text = str(col_name)
#                 style_header_cell(cell)

#             # Data rows
#             for row_idx, row_data in enumerate(
#                 df.iloc[start_row:start_row + rows_per_slide].values, start=1
#             ):
#                 for col_idx, value in enumerate(row_data):
#                     style_data_cell(table.cell(row_idx, col_idx), value)
                    
#                 # ✅ adjust height after filling row
#                 adjust_row_height(table, row_idx, row_data, col_widths)

#             start_row += rows_per_slide
#             slide_count += 1

#     prs.save(output_path)
#     return output_path


# def create_ppt_from_excel(template_path, excel_data, output_path, rows_per_slide=5):
#     """
#     Creates a PPT from Excel sheet data (dict of {sheet_name: dataframe}).
#     Keeps the first TWO slides of the template as-is, then inserts generated slides,
#     and leaves the remaining template slides at the end.
#     """
#     prs = Presentation(template_path)

#     # insertion starts after the 2nd slide (0-based index = 1), so position = 2
#     base_insert_index = 2
#     added_so_far = 0  # how many of our generated slides we’ve inserted

#     for sheet_name, df in excel_data.items():
#         if df.empty or len(df.columns) == 0:
#             continue

#         total_rows = len(df)
#         start_row = 0
#         slide_count = 1

#         while start_row < total_rows:
#             slide_layout = prs.slide_layouts[1]  # Title & Content (adjust if your template uses another)
#             # insert at (2 + already_inserted_count)
#             slide = insert_slide(prs, slide_layout, base_insert_index + added_so_far)
#             added_so_far += 1

#             # Title
#             title_text = f"{sheet_name} (Contd..)" if slide_count > 1 else sheet_name
#             if slide.shapes.title:
#                 slide.shapes.title.text = title_text
#             else:
#                 left, top, width, height = Inches(0.5), Inches(0.3), Inches(9), Inches(1)
#                 textbox = slide.shapes.add_textbox(left, top, width, height)
#                 textbox.text_frame.text = title_text

#             # Table
#             rows = min(rows_per_slide, total_rows - start_row) + 1  # +1 for header
#             cols = len(df.columns)
#             left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.8 * rows)

#             table = slide.shapes.add_table(rows, cols, left, top, width, height).table

#             # Column widths (even split across a target width; tweak if needed)
#             total_table_width = Inches(12)
#             num_cols = len(df.columns)
#             col_widths = [Inches(total_table_width.inches / num_cols)] * num_cols
#             for i, w in enumerate(col_widths[:cols]):
#                 table.columns[i].width = int(w.emu)

#             # Header row
#             for col_idx, col_name in enumerate(df.columns):
#                 cell = table.cell(0, col_idx)
#                 cell.text = str(col_name)
#                 style_header_cell(cell)

#             # Data rows
#             for row_idx, row_data in enumerate(
#                 df.iloc[start_row:start_row + rows_per_slide].values, start=1
#             ):
#                 for col_idx, value in enumerate(row_data):
#                     style_data_cell(table.cell(row_idx, col_idx), value)
#                 adjust_row_height(table, row_idx, row_data, col_widths)

#             start_row += rows_per_slide
#             slide_count += 1

#     prs.save(output_path)
#     return output_path

#update project title in slide 1
def update_project_title(slide, new_title):
    for shape in slide.shapes:
        if shape.name == 'Rectangle 6' and shape.has_text_frame:
            paragraphs = shape.text_frame.paragraphs
            if len(paragraphs) > 0:
                paragraph = paragraphs[0]
                # Save formatting from first run (font, color, size etc)
                if paragraph.runs:
                    orig_run = paragraph.runs[0]
                    orig_font = orig_run.font
                    font_dict = {
                        'size': orig_font.size,
                        'bold': orig_font.bold,
                        'italic': orig_font.italic,
                        'underline': orig_font.underline,
                        'color': orig_font.color.rgb if orig_font.color and orig_font.color.rgb else None,
                    }
                else:
                    font_dict = {}

                # Remove all runs
                while paragraph.runs:
                    paragraph._p.remove(paragraph.runs[0]._r)
                # Add new run/text
                run = paragraph.add_run()
                run.text = new_title
                
                # Force font family to Calibri Bold
                run.font.name = "Calibri"
                run.font.bold = True

                # Reapply original formatting
                if font_dict.get('size'): run.font.size = font_dict['size']
                if font_dict.get('bold') is not None: run.font.bold = font_dict['bold']
                if font_dict.get('italic') is not None: run.font.italic = font_dict['italic']
                if font_dict.get('underline') is not None: run.font.underline = font_dict['underline']
                if font_dict.get('color'): run.font.color.rgb = font_dict['color']
                # Optionally also set alignment if you want:
                # paragraph.alignment = PP_ALIGN.CENTER  # etc, as per your template
                return
            
#Calculate this week's Friday date
def get_next_friday():
    today = datetime.date.today()
    weekday = today.weekday()  # Monday=0 ... Sunday=6

    # Calculate days until next Friday
    days_ahead = (4 - weekday) % 7  # Friday = 4
    if days_ahead == 0:
        days_ahead = 7  # If today is Friday, pick NEXT Friday

    next_friday = today + datetime.timedelta(days=days_ahead)
    return next_friday.strftime("%d %B %Y")  # Example: "12 September 2025"



#Upadate Friday date in slide
def update_slide_date(slide):
    friday_str = get_next_friday()

    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if "Date:" in p.text:   # Look for the Date line
                    # Save formatting
                    if p.runs:
                        orig_run = p.runs[0]
                        font_dict = {
                            'size': orig_run.font.size,
                            'bold': orig_run.font.bold,
                            'italic': orig_run.font.italic,
                            'underline': orig_run.font.underline,
                            'color': orig_run.font.color.rgb if orig_run.font.color and orig_run.font.color.rgb else None,
                        }
                    else:
                        font_dict = {}

                    # Remove old runs
                    while p.runs:
                        p._p.remove(p.runs[0]._r)

                    # Add new text
                    run = p.add_run()
                    run.text = f"Date:  {friday_str}"
                    
                    # Force font family to Calibri Bold
                    run.font.name = "Calibri"
                    run.font.bold = True

                    # Reapply formatting
                    if font_dict.get('size'): run.font.size = font_dict['size']
                    if font_dict.get('bold') is not None: run.font.bold = font_dict['bold']
                    if font_dict.get('italic') is not None: run.font.italic = font_dict['italic']
                    if font_dict.get('underline') is not None: run.font.underline = font_dict['underline']
                    if font_dict.get('color'): run.font.color.rgb = font_dict['color']
                    return

def copy_textbox_to_slide(source_slide, target_slide, shape_name="TextBox 1"):
    """
    Copy a specific text box (by name) from source_slide to target_slide.
    Keeps it editable.
    """
    for shape in source_slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            el = deepcopy(shape.element)  # clone XML element
            target_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
            return


def create_ppt_from_excel(template_path, excel_data, output_path, rows_per_slide=5, project_title=""):
    """
    Creates a PPT from Excel sheet data.
    Supports both:
      - pandas.DataFrame
      - list of row dicts with nested {"value": ..., "style": ...}
    """
    prs = Presentation(template_path)
    if project_title:
      update_project_title(prs.slides[0], project_title)
      
    # ✅ Update Friday date (assumes your template has a shape named "Date Placeholder")
    update_slide_date(prs.slides[0])
      
    print("=== First Slide Paragraphs for Debug ===")
    for shape in prs.slides[0].shapes:
       if shape.has_text_frame:
          for i, p in enumerate(shape.text_frame.paragraphs):
             print(f"Shape: '{shape.name}', Paragraph {i}: {p.text}")
    print("=== End Debug ===")


    base_insert_index = 2
    added_so_far = 0

    for sheet_name, sheet_content in excel_data.items():
        # --- Detect if it's a DataFrame or list ---
        if isinstance(sheet_content, pd.DataFrame):
            if sheet_content.empty or len(sheet_content.columns) == 0:
                continue
            data_rows = sheet_content.to_dict(orient="records")
            columns = list(sheet_content.columns)
        elif isinstance(sheet_content, list) and len(sheet_content) > 0:
            # JSON rows (list of dicts with value+style)
            data_rows = sheet_content
            columns = list(sheet_content[0].keys())
        else:
            continue

        total_rows = len(data_rows)
        start_row = 0
        slide_count = 1

        while start_row < total_rows:
            slide_layout = prs.slide_layouts[1]  
            slide = insert_slide(prs, slide_layout, base_insert_index + added_so_far)
            added_so_far += 1

            title_text = f"{sheet_name} (Contd..)" if slide_count > 1 else sheet_name
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            else:
                left, top, width, height = Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                textbox.text_frame.text = title_text

            # Table setup
            rows = min(rows_per_slide, total_rows - start_row) + 1  
            cols = len(columns)
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.8 * rows)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Column widths
            total_table_width = Inches(12)
            col_widths = [Inches(total_table_width.inches / cols)] * cols
            for i, w in enumerate(col_widths):
                table.columns[i].width = int(w.emu)

            # Header row
            for col_idx, col_name in enumerate(columns):
                cell = table.cell(0, col_idx)
                cell.text = str(col_name)
                style_header_cell(cell)

            # Data rows
            chunk = data_rows[start_row:start_row + rows_per_slide]
            for row_idx, row_data in enumerate(chunk, start=1):
                for col_idx, col_name in enumerate(columns):
                    cell_info = row_data.get(col_name, "")
                    style_data_cell(table.cell(row_idx, col_idx), cell_info)
                adjust_row_height(
                    table, row_idx,
                    [row_data.get(c, "") for c in columns],
                    col_widths
                )

            start_row += rows_per_slide
            slide_count += 1

    prs.save(output_path)
    return output_path
